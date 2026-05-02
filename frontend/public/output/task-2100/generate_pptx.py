#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
和光智成融资路演PPT生成脚本
Pre-A轮融资 | 2026年Q2
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_pitch_deck():
    prs = Presentation()
    
    # 设置默认字体
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 幻灯片1: 封面页
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 50, 100)
    
    # 主标题
    left = Inches(1)
    top = Inches(2)
    width = Inches(11)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "和光智成 Helight"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    left = Inches(1)
    top = Inches(3.5)
    width = Inches(11)
    height = Inches(1)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = "北航背书的AI材料+药物双擎发现平台"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(200, 220, 255)
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = tf2.add_paragraph()
    p3.text = "—— 中国的Periodic Labs"
    p3.font.size = Pt(24)
    p3.font.color.rgb = RGBColor(180, 200, 230)
    p3.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    left = Inches(1)
    top = Inches(5.5)
    width = Inches(11)
    height = Inches(1)
    txBox3 = slide.shapes.add_textbox(left, top, width, height)
    tf3 = txBox3.text_frame
    p4 = tf3.add_paragraph()
    p4.text = "Pre-A轮融资  |  5,000万-1亿 RMB  |  目标估值 5-8亿"
    p4.font.size = Pt(22)
    p4.font.color.rgb = RGBColor(255, 255, 255)
    p4.alignment = PP_ALIGN.CENTER
    
    p5 = tf3.add_paragraph()
    p5.text = "2026年Q2"
    p5.font.size = Pt(18)
    p5.font.color.rgb = RGBColor(200, 200, 200)
    p5.alignment = PP_ALIGN.CENTER
    
    # 幻灯片2: 执行摘要
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "执行摘要 (Executive Summary)"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(25, 50, 100)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "【投资亮点】"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    
    points = [
        "赛道爆发：生成式AI材料市场2034年预计$117亿，中国CAGR 31.2%全球领先",
        "对标验证：Periodic Labs 12个月估值从$1.3B→$70亿（翻7倍），验证商业模式",
        "技术壁垒：北航联合实验室 + AI算法 + 湿实验自动化闭环",
        "团队顶配：北航教授领衔，学术界+产业界复合型团队",
        "里程碑明确：2027年实现盈亏平衡，2028年启动IPO准备"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【融资需求】"
    p.font.bold = True
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "融资金额：5,000万-10,000万 RMB  |  出让股权：10%-15%  |  目标估值：5-8亿 RMB"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "资金用途：研发40% | 市场30% | 团队20% | 运营10%"
    p.level = 1
    p.font.size = Pt(18)
    
    # 幻灯片3: 市场机会
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "市场机会 - 万亿级蓝海"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(25, 50, 100)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "【全球AI材料科学市场】"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    
    market_data = [
        "2026年全球AI融资总额：> $1,300亿",
        "AI占全球VC比例：61%",
        "生成式AI材料市场(2034E)：$117亿",
        "中国AI材料市场CAGR：31.2%"
    ]
    
    for data in market_data:
        p = tf.add_paragraph()
        p.text = data
        p.level = 1
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【政策红利】"
    p.font.bold = True
    p.font.size = Pt(20)
    
    policies = [
        "北京\"AI+新材料\"战略全面推进",
        "房山\"首都材料谷\"180亿产值集群",
        "中试放大平台政策补贴支持",
        "Stanford 2026 AI Index：中国论文/专利/工业机器人安装量全球第一"
    ]
    
    for policy in policies:
        p = tf.add_paragraph()
        p.text = policy
        p.level = 1
        p.font.size = Pt(18)
    
    # 幻灯片4: 竞品对标
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "竞品对标分析"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(25, 50, 100)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "【全球竞品估值与营收对比】"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    
    competitors = [
        "Periodic Labs：$70亿估值 (2026.3)，12个月翻7倍",
        "XtalPi (晶泰科技)：2025年营收8.026亿 RMB，+201% YoY",
        "Schrödinger：$2.56亿营收，软件收入+$1.995亿",
        "Insilico Medicine：$30亿估值，AI药物研发领先",
        "Kebotix：$12亿估值，AI材料发现"
    ]
    
    for comp in competitors:
        p = tf.add_paragraph()
        p.text = comp
        p.level = 1
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【和光智成差异化优势】"
    p.font.bold = True
    p.font.size = Pt(20)
    
    advantages = [
        "北航背书：空天动力结构强度智能实验北京市重点实验室共建单位",
        "双轮驱动：材料发现 + 药物发现，分散风险",
        "中国成本：同等能力下成本仅为海外竞品1/3",
        "闭环实验：AI算法 + 自动化湿实验平台"
    ]
    
    for adv in advantages:
        p = tf.add_paragraph()
        p.text = adv
        p.level = 1
        p.font.size = Pt(18)
    
    # 幻灯片5: 技术平台
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "技术平台 - AI+实验闭环"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(25, 50, 100)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "【Helight AI材料发现平台】"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "全流程AI赋能：分子设计 → 性能预测 → 合成规划 → 自动实验 → 数据反馈"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【核心技术能力】"
    p.font.bold = True
    p.font.size = Pt(20)
    
    techs = [
        "分子生成模型：基于Transformer的逆合成预测，准确率>85%",
        "性质预测引擎：100+材料性质预测模型，MAE<5%",
        "实验自动化：集成液体处理机器人，通量提升10倍",
        "知识图谱：1亿+化合物，1000万+文献数据"
    ]
    
    for tech in techs:
        p = tf.add_paragraph()
        p.text = tech
        p.level = 1
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【技术护城河：数据壁垒 + 算法壁垒 + 实验壁垒 + 人才壁垒】"
    p.font.bold = True
    p.font.size = Pt(18)
    
    # 幻灯片6: 商业模式
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "商业模式"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(25, 50, 100)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "【收入结构演变】"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    
    revenue_struct = [
        "2026年：技术服务80% + 项目合作20%",
        "2027年：技术服务50% + SaaS 30% + 项目合作20%",
        "2028年：SaaS 30% + 技术服务40% + 项目合作30%"
    ]
    
    for rs in revenue_struct:
        p = tf.add_paragraph()
        p.text = rs
        p.level = 1
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【定价策略】"
    p.font.bold = True
    p.font.size = Pt(20)
    
    prices = [
        "Tier 1 (头部企业)：200-500万/年，3-5家，贡献60%收入",
        "Tier 2 (中型企业)：50-200万/年，10-15家，贡献30%收入",
        "Tier 3 (小型企业)：10-50万/年，20-30家，贡献8%收入",
        "Tier 4 (高校/科研)：1-10万/年，50+家，贡献2%收入"
    ]
    
    for price in prices:
        p = tf.add_paragraph()
        p.text = price
        p.level = 1
        p.font.size = Pt(18)
    
    # 幻灯片7: 最新业绩与里程碑
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "2026 Q2 最新业绩与里程碑"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(25, 50, 100)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "【已达成里程碑 ✅】"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    
    achieved = [
        "实验室资质：空天动力结构强度智能实验北京市重点实验室共建单位，刘宇宙教授任副主任",
        "客户突破：签约首批3家付费客户，在谈意向客户10+家",
        "融资进展：陈俊宇（陈总）主动索要BP，水木创投李林泽主动联系"
    ]
    
    for a in achieved:
        p = tf.add_paragraph()
        p.text = a
        p.level = 1
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【Q2-Q3 目标 🎯】"
    p.font.bold = True
    p.font.size = Pt(20)
    
    goals = [
        "完成Pre-A轮融资 (5000万-1亿)",
        "签约10家付费客户，营收突破500万",
        "申请发明专利5-8项"
    ]
    
    for g in goals:
        p = tf.add_paragraph()
        p.text = g
        p.level = 1
        p.font.size = Pt(18)
    
    # 幻灯片8: 财务预测
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "财务预测 (2026-2029)"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(25, 50, 100)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "【收入预测 (单位：万元 RMB)】"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    
    financials = [
        "2026E：总营收 600万，毛利率65%，净利润-170万",
        "2027E：总营收 3,000万，毛利率72%，净利润552万 (+400% YoY)",
        "2028E：总营收 8,000万，毛利率78%，净利润2,000万 (+167% YoY)",
        "2029E：总营收 20,000万，毛利率82%，净利润6,000万 (+150% YoY)"
    ]
    
    for f in financials:
        p = tf.add_paragraph()
        p.text = f
        p.level = 1
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【关键财务指标】"
    p.font.bold = True
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "2027年实现盈亏平衡 | 2028年净利润率达25% | 2029年启动IPO准备"
    p.level = 1
    p.font.size = Pt(18)
    
    # 幻灯片9: 估值测算
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "估值测算"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(25, 50, 100)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "【方法一：可比公司法】"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "可比公司估值/营收倍数：14x-58x，平均37.8x"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "和光智成2027E前瞻估值：3.6亿 - 4.5亿 RMB (12x-15x)"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【方法二：DCF现金流折现法】"
    p.font.bold = True
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "关键假设：WACC 25%，终值增长率5%"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "DCF计算结果：企业价值约6.11亿，股权价值约6.86亿 RMB"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【综合估值结论：5 - 8亿 RMB】"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(200, 50, 50)
    
    # 幻灯片10: 融资计划与资金用途
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "融资计划与资金用途"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(25, 50, 100)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "【本轮融资】"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "轮次：Pre-A轮 | 金额：5,000万-10,000万 RMB | 出让股权：10%-15% | 目标估值：5-8亿 RMB"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【资金用途】"
    p.font.bold = True
    p.font.size = Pt(20)
    
    uses = [
        "研发投入 40% (3,000万)：AI算法研发、实验平台建设、专利布局",
        "市场拓展 30% (2,250万)：销售团队、市场推广、客户成功",
        "团队建设 20% (1,500万)：核心人才招聘、股权激励",
        "运营资金 10% (750万)"
    ]
    
    for use in uses:
        p = tf.add_paragraph()
        p.text = use
        p.level = 1
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【投资人优先级：北航系基金 → 硬科技VC → 产业资本 → 政府引导基金 → 天使投资人】"
    p.font.bold = True
    p.font.size = Pt(16)
    
    # 幻灯片11: 里程碑与退出路径
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "里程碑与退出路径"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(25, 50, 100)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "【关键里程碑】"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    
    milestones = [
        "2026 Q3：完成Pre-A轮融资",
        "2026 Q4：签约10+客户，营收500万+",
        "2027 Q2：实现盈亏平衡",
        "2027 Q4：年营收3,000万，启动A轮",
        "2028 Q4：年营收8,000万，净利润2,000万",
        "2029 H2：启动IPO准备工作"
    ]
    
    for m in milestones:
        p = tf.add_paragraph()
        p.text = m
        p.level = 1
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【退出路径】"
    p.font.bold = True
    p.font.size = Pt(20)
    
    exits = [
        "IPO：科创板/北交所（2029-2030）",
        "并购：头部材料企业/药企战略收购",
        "下一轮退出：A/B轮老股转让"
    ]
    
    for e in exits:
        p = tf.add_paragraph()
        p.text = e
        p.level = 1
        p.font.size = Pt(18)
    
    # 幻灯片12: 联系方式
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 50, 100)
    
    left = Inches(1)
    top = Inches(2)
    width = Inches(11)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "感谢您的关注！"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    left = Inches(1)
    top = Inches(3.5)
    width = Inches(11)
    height = Inches(2)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = "和光智成（北京）科技有限公司"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(200, 220, 255)
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = tf2.add_paragraph()
    p3.text = "地址：北京市海淀区北京航空航天大学"
    p3.font.size = Pt(20)
    p3.font.color.rgb = RGBColor(180, 200, 230)
    p3.alignment = PP_ALIGN.CENTER
    
    p4 = tf2.add_paragraph()
    p4.text = "联系人：刘宇宙 教授"
    p4.font.size = Pt(20)
    p4.font.color.rgb = RGBColor(180, 200, 230)
    p4.alignment = PP_ALIGN.CENTER
    
    # 保存PPT
    output_path = "/Users/mettlyz/.openclaw/workspace/output/task-2100/和光智成_融资路演PitchDeck_Q2_20260426.pptx"
    prs.save(output_path)
    print(f"PPT已生成：{output_path}")
    print(f"共 {len(prs.slides)} 张幻灯片")

if __name__ == "__main__":
    create_pitch_deck()
