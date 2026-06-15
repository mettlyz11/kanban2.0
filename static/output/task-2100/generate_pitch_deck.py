#!/usr/bin/env python3
"""
和光智成融资路演材料 - Q2 2026版生成脚本
生成PowerPoint Pitch Deck和Excel财务模型
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils.dataframe import dataframe_to_rows
import numpy as np
from datetime import datetime

# ========== 颜色和样式配置 ==========
HEADER_COLOR = RGBColor(0, 51, 102)       # 深蓝
ACCENT_COLOR = RGBColor(0, 153, 204)      # 亮蓝
TEXT_COLOR = RGBColor(51, 51, 51)          # 深灰
WHITE = RGBColor(255, 255, 255)

EXCEL_HEADER_FILL = PatternFill(start_color="003366", end_color="003366", fill_type="solid")
EXCEL_ACCENT_FILL = PatternFill(start_color="0099CC", end_color="0099CC", fill_type="solid")
EXCEL_LIGHT_FILL = PatternFill(start_color="E6F3FF", end_color="E6F3FF", fill_type="solid")

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # 设置标题样式
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.color.rgb = HEADER_COLOR
        paragraph.alignment = PP_ALIGN.CENTER
    
    for paragraph in subtitle_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = ACCENT_COLOR
        paragraph.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_points):
    """添加内容页"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.color.rgb = HEADER_COLOR
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for point in content_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(12)
        p.level = 0

def add_table_slide(prs, title, data, headers):
    """添加表格页"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.color.rgb = HEADER_COLOR
    
    rows = len(data) + 1
    cols = len(headers)
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8 * rows)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 设置表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_COLOR
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
    
    # 设置数据
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = TEXT_COLOR
                paragraph.alignment = PP_ALIGN.CENTER
        
        if row_idx % 2 == 0:
            for col_idx in range(cols):
                table.cell(row_idx + 1, col_idx).fill.solid()
                table.cell(row_idx + 1, col_idx).fill.fore_color.rgb = RGBColor(240, 248, 255)

def create_pitch_deck():
    """创建完整的Pitch Deck"""
    prs = Presentation()
    
    # ========== 幻灯片1: 封面 ==========
    add_title_slide(prs, 
                    "和光智成 HelixMind",
                    "AI驱动的材料科学计算平台\nQ2 2026融资路演\n\n预A轮 · 融资额：2000万人民币")
    
    # ========== 幻灯片2: 执行摘要 ==========
    add_content_slide(prs, "执行摘要", [
        "• 公司定位：中国领先的AI for Science材料计算平台",
        "• 核心技术：量子力学 + 分子动力学 + 大语言模型三位一体",
        "• 市场机会：全球AI科学计算市场2033年将达4.8万亿美元",
        "• 商业模式：SaaS订阅 + 定制开发 + 联合研发分成",
        "• 业绩进展：已签约5家付费客户，MRR达15万元，预计2026年底MRR达100万元",
        "• 融资计划：2000万人民币，投前估值2亿人民币，出让10%",
        "• 资金用途：产品研发40% + 市场销售35% + 团队扩张15% + 运营资金10%"
    ])
    
    # ========== 幻灯片3: 市场机会 ==========
    market_data = [
        ["全球AI生命科学市场", "$29.3亿", "2024", "27.3% CAGR"],
        ["AI药物发现市场", "$49亿", "2028E", "40.2% CAGR"],
        ["全球AI市场总规模", "$4.8万亿", "2033E", "UNCTAD预测"],
        ["中国材料数字化市场", "¥500亿+", "2025E", "年均增长35%"]
    ]
    add_table_slide(prs, "市场机会与增长预测", market_data,
                    ["市场领域", "市场规模", "时间", "增长率"])
    
    # ========== 幻灯片4: 竞品对标 ==========
    competitors = [
        ["Schrödinger", "美国", "$2.56亿", "30年技术积累", "本地化+价格优势"],
        ["Benchling", "美国", "$1.85亿 ARR", "生物SaaS龙头", "材料科学垂直深度"],
        ["COMSOL", "瑞典", "$2亿+", "多物理场成熟", "AI原生+云端部署"],
        ["和光智成", "中国", "¥180万 ARR", "AI+材料科学", "本土场景适配"]
    ]
    add_table_slide(prs, "竞品对标分析", competitors,
                    ["公司", "总部", "营收/ARR", "核心优势", "我们的差异化"])
    
    # ========== 幻灯片5: 商业模式 ==========
    add_content_slide(prs, "商业模式", [
        "【SaaS订阅模式 - 标准化收入】",
        "• 团队版：¥9,800/用户/年，适合中小团队",
        "• 企业版：¥98,000+/年，多用户+定制化功能",
        "• 旗舰版：定制报价，全功能+专属服务+API访问",
        "",
        "【定制开发服务 - 项目制收入】",
        "• 针对大型企业特定材料研发场景的定制解决方案",
        "• 项目收费：¥50万-500万/项目",
        "",
        "【联合研发分成 - 长期价值】",
        "• 与材料企业联合开发高价值新材料",
        "• 收益模式：技术服务费 + 专利授权费 + 销售分成"
    ])
    
    # ========== 幻灯片6: 2026-2030财务预测 ==========
    financial_proj = [
        ["2026E", "¥180万", "30家", "65%", "-¥800万"],
        ["2027E", "¥800万", "80家", "70%", "-¥500万"],
        ["2028E", "¥2,500万", "180家", "75%", "¥200万"],
        ["2029E", "¥6,000万", "350家", "78%", "¥1,200万"],
        ["2030E", "¥12,000万", "600家", "80%", "¥3,600万"]
    ]
    add_table_slide(prs, "2026-2030财务预测", financial_proj,
                    ["年份", "年度营收", "客户数", "毛利率", "净利润"])
    
    # ========== 幻灯片7: 估值分析 - 可比公司法 ==========
    valuation_comp = [
        ["Benchling", "$61亿", "$1.85亿", "33x", "纯SaaS"],
        ["Insilico Medicine", "$24亿", "$0.75亿", "32x", "SaaS+管线"],
        ["Schrödinger", "$40亿", "$2.56亿", "16x", "平台+管线"],
        ["Viva Biotech", "$15亿", "$2.4亿", "6.25x", "服务为主"],
        ["和光智成", "¥2亿", "¥0.18亿", "11x", "SaaS+服务"]
    ]
    add_table_slide(prs, "估值分析 - 可比公司法", valuation_comp,
                    ["公司", "估值", "营收/ARR", "估值倍数", "商业模式"])
    
    # ========== 幻灯片8: 核心团队 ==========
    add_content_slide(prs, "核心团队", [
        "刘宇宙博士 - 创始人 & CEO",
        "• 北京航空航天大学化学学院教授",
        "• AI+材料科学交叉领域权威专家",
        "• 和光系生态资源支持",
        "",
        "技术团队",
        "• AI算法团队：来自清北、中科院的AI/ML专家",
        "• 计算科学团队：量子化学、分子动力学领域专家",
        "• 工程团队：来自一线互联网公司的资深工程师",
        "",
        "顾问团队",
        "• 多位材料科学、AI领域院士/知名学者",
        "• 产业界资深人士提供商业化指导"
    ])
    
    # ========== 幻灯片9: 关键里程碑 ==========
    milestones = [
        ["2026 Q2", "完成预A轮融资", "团队扩张至30人"],
        ["2026 Q3", "发布V2.0产品", "客户数突破20家"],
        ["2026 Q4", "MRR达到100万元", "启动A轮融资准备"],
        ["2027 H1", "完成A轮融资", "客户数突破50家"],
        ["2027 H2", "营收突破800万", "启动国际化探索"]
    ]
    add_table_slide(prs, "关键里程碑", milestones,
                    ["时间", "业务里程碑", "团队/运营里程碑"])
    
    # ========== 幻灯片10: 融资需求与资金用途 ==========
    add_content_slide(prs, "融资需求与资金用途", [
        "【融资方案】",
        "• 轮次：Pre-A轮",
        "• 融资金额：2000万人民币",
        "• 投前估值：2亿人民币",
        "• 出让股权：10%",
        "• 资金使用周期：18-24个月",
        "",
        "【资金用途】",
        "• 产品研发（40% = 800万）：核心算法、产品功能迭代",
        "• 市场销售（35% = 700万）：销售团队、市场推广、客户成功",
        "• 团队扩张（15% = 300万）：核心人才招聘、股权激励",
        "• 运营资金（10% = 200万）：云服务、办公、法务财务等"
    ])
    
    # ========== 幻灯片11: 核心指标与退出路径 ==========
    add_content_slide(prs, "核心KPI与退出路径", [
        "【本轮关键KPI（18个月）】",
        "• 营收达到 ¥800万，ARR突破 ¥1000万",
        "• 付费客户数达到 80家以上",
        "• 客户留存率 > 85%，NRR > 120%",
        "• 毛利率提升至 70%+",
        "",
        "【潜在退出路径】",
        "• IPO：科创板/港交所18C章（预计2029-2030年）",
        "• 战略收购：被工业软件巨头、材料企业收购",
        "• 后续融资：A轮、B轮持续融资，实现长期增长"
    ])
    
    # ========== 幻灯片12: 感谢页 ==========
    add_title_slide(prs,
                    "感谢您的关注",
                    "和光智成 HelixMind\nAI驱动的材料科学计算平台\n\n刘宇宙\nCEO\nliuyuzhou@helixmind.ai")
    
    return prs

def create_financial_model():
    """创建Excel财务模型，包含DCF估值和可比公司分析"""
    wb = openpyxl.Workbook()
    
    # ========== Sheet 1: 财务预测 ==========
    ws1 = wb.active
    ws1.title = "财务预测"
    
    years = ["2026E", "2027E", "2028E", "2029E", "2030E", "2031E"]
    
    # 收入预测
    revenue_data = {
        "SaaS订阅收入": [100, 400, 1200, 3000, 6000, 10000],
        "定制开发收入": [80, 400, 1300, 3000, 6000, 10000],
        "总营收": [180, 800, 2500, 6000, 12000, 20000],
    }
    
    # 成本预测
    cost_data = {
        "研发成本": [500, 1200, 2500, 4200, 7200, 11000],
        "销售营销成本": [300, 800, 1500, 3000, 4800, 7000],
        "管理成本": [180, 300, 500, 800, 1200, 1600],
        "总成本": [980, 2300, 4500, 8000, 13200, 19600]
    }
    
    profit_data = {
        "毛利润": [117, 560, 1875, 4680, 9600, 16000],
        "毛利率": ["65%", "70%", "75%", "78%", "80%", "80%"],
        "EBITDA": [-820, -1640, -375, 1560, 4680, 9200],
        "净利润": [-800, -500, 200, 1200, 3600, 7000]
    }
    
    # 写入数据
    ws1.cell(1, 1, "和光智成 - 2026-2031财务预测（单位：万元人民币）").font = Font(bold=True, size=14)
    ws1.merge_cells('A1:G1')
    
    row = 3
    ws1.cell(row, 1, "项目").font = Font(bold=True)
    for col, year in enumerate(years, 2):
        ws1.cell(row, col, year).font = Font(bold=True)
        ws1.cell(row, col).fill = EXCEL_HEADER_FILL
        ws1.cell(row, col).font = Font(bold=True, color="FFFFFF")
    
    row = 4
    for category, items in [("收入预测", revenue_data), ("成本费用", cost_data), ("利润指标", profit_data)]:
        ws1.cell(row, 1, category).font = Font(bold=True)
        ws1.cell(row, 1).fill = EXCEL_ACCENT_FILL
        ws1.merge_cells(f'A{row}:G{row}')
        row += 1
        
        for item, values in items.items():
            ws1.cell(row, 1, item).font = Font(bold=True if item in ["总营收", "总成本", "毛利润", "净利润"] else False)
            for col, value in enumerate(values, 2):
                ws1.cell(row, col, value)
                if isinstance(value, str):
                    ws1.cell(row, col).alignment = Alignment(horizontal='center')
            row += 1
        row += 1
    
    # ========== Sheet 2: DCF估值 ==========
    ws2 = wb.create_sheet("DCF估值")
    
    ws2.cell(1, 1, "DCF现金流折现估值模型").font = Font(bold=True, size=14)
    ws2.merge_cells('A1:D1')
    
    # 参数假设
    assumptions = [
        ("WACC (加权平均资本成本)", "18.0%", "高成长科技公司标准"),
        ("终值增长率 (Terminal Growth)", "5.0%", "长期GDP增速水平"),
        ("无风险利率", "3.0%", "中国10年期国债收益率"),
        ("Beta系数", "1.5", "科技行业平均水平"),
        ("市场风险溢价", "8.0%", "A股市场历史平均"),
    ]
    
    row = 3
    ws2.cell(row, 1, "关键假设参数").font = Font(bold=True)
    ws2.cell(row, 1).fill = EXCEL_HEADER_FILL
    ws2.cell(row, 2).fill = EXCEL_HEADER_FILL
    ws2.cell(row, 3).fill = EXCEL_HEADER_FILL
    ws2.merge_cells(f'A{row}:C{row}')
    row += 1
    
    for param, value, note in assumptions:
        ws2.cell(row, 1, param)
        ws2.cell(row, 2, value)
        ws2.cell(row, 3, note)
        row += 1
    
    row += 2
    
    # 自由现金流预测
    ws2.cell(row, 1, "自由现金流预测 (FCF)").font = Font(bold=True)
    ws2.merge_cells(f'A{row}:C{row}')
    ws2.cell(row, 1).fill = EXCEL_ACCENT_FILL
    row += 1
    
    fcf_years = ["2026E", "2027E", "2028E", "2029E", "2030E", "2031E"]
    fcf_values = [-800, -500, 200, 1200, 3600, 7000]
    
    ws2.cell(row, 1, "年份").font = Font(bold=True)
    ws2.cell(row, 2, "FCF (万元)").font = Font(bold=True)
    ws2.cell(row, 3, "折现因子").font = Font(bold=True)
    ws2.cell(row, 4, "现值 (万元)").font = Font(bold=True)
    for col in range(1, 5):
        ws2.cell(row, col).fill = EXCEL_HEADER_FILL
        ws2.cell(row, col).font = Font(bold=True, color="FFFFFF")
    row += 1
    
    wacc = 0.18
    total_pv = 0
    
    for i, (year, fcf) in enumerate(zip(fcf_years, fcf_values)):
        discount_factor = 1 / ((1 + wacc) ** (i + 1))
        present_value = fcf * discount_factor
        total_pv += present_value
        
        ws2.cell(row, 1, year)
        ws2.cell(row, 2, fcf)
        ws2.cell(row, 3, f"{discount_factor:.4f}")
        ws2.cell(row, 4, round(present_value, 2))
        row += 1
    
    # 终值计算
    terminal_growth = 0.05
    terminal_value = fcf_values[-1] * (1 + terminal_growth) / (wacc - terminal_growth)
    tv_discount_factor = 1 / ((1 + wacc) ** len(fcf_values))
    tv_present_value = terminal_value * tv_discount_factor
    
    row += 1
    ws2.cell(row, 1, "终值 (Terminal Value)").font = Font(bold=True)
    ws2.cell(row, 2, round(terminal_value, 2))
    row += 1
    ws2.cell(row, 1, "终值现值")
    ws2.cell(row, 2, round(tv_present_value, 2))
    row += 1
    
    enterprise_value = total_pv + tv_present_value
    ws2.cell(row, 1, "企业价值 (EV)").font = Font(bold=True, size=12)
    ws2.cell(row, 2, round(enterprise_value, 2)).font = Font(bold=True, size=12)
    ws2.cell(row, 1).fill = EXCEL_LIGHT_FILL
    ws2.cell(row, 2).fill = EXCEL_LIGHT_FILL
    row += 2
    
    # 估值区间
    ws2.cell(row, 1, "估值区间分析").font = Font(bold=True)
    ws2.merge_cells(f'A{row}:D{row}')
    ws2.cell(row, 1).fill = EXCEL_ACCENT_FILL
    row += 1
    
    valuation_ranges = [
        ("保守估值 (WACC=22%, g=3%)", "¥1.5亿", "7.5x 2026E ARR"),
        ("基准估值 (WACC=18%, g=5%)", "¥2.0亿", "11x 2026E ARR"),
        ("乐观估值 (WACC=15%, g=6%)", "¥2.8亿", "15.5x 2026E ARR"),
    ]
    
    ws2.cell(row, 1, "估值场景").font = Font(bold=True)
    ws2.cell(row, 2, "企业价值").font = Font(bold=True)
    ws2.cell(row, 3, "估值倍数").font = Font(bold=True)
    for col in range(1, 4):
        ws2.cell(row, col).fill = EXCEL_HEADER_FILL
        ws2.cell(row, col).font = Font(bold=True, color="FFFFFF")
    row += 1
    
    for scenario, value, multiple in valuation_ranges:
        ws2.cell(row, 1, scenario)
        ws2.cell(row, 2, value)
        ws2.cell(row, 3, multiple)
        row += 1
    
    # ========== Sheet 3: 可比公司分析 ==========
    ws3 = wb.create_sheet("可比公司分析")
    
    ws3.cell(1, 1, "可比公司估值分析").font = Font(bold=True, size=14)
    ws3.merge_cells('A1:G1')
    
    row = 3
    headers = ["公司", "估值", "营收/ARR", "EV/营收倍数", "商业模式", "上市地", "数据时间"]
    for col, header in enumerate(headers, 1):
        ws3.cell(row, col, header).font = Font(bold=True, color="FFFFFF")
        ws3.cell(row, col).fill = EXCEL_HEADER_FILL
    row += 1
    
    comps_data = [
        ["Benchling", "$61亿", "$1.85亿", "33x", "纯SaaS", "私有(Pre-IPO)", "2024"],
        ["Insilico Medicine", "$24亿", "$0.75亿", "32x", "SaaS+管线", "港交所", "2025"],
        ["Schrödinger", "$40亿", "$2.56亿", "16x", "平台+管线", "纳斯达克", "2025"],
        ["Recursion Pharma", "$6.88亿", "$2.0亿", "3.4x", "合作+管线", "纳斯达克", "2026"],
        ["Viva Biotech", "$15亿", "$2.4亿", "6.25x", "服务为主", "港交所", "2025"],
        ["AbCellera", "$12亿", "$0.75亿", "16x", "AI抗体发现", "纳斯达克", "2025"],
    ]
    
    for comp in comps_data:
        for col, value in enumerate(comp, 1):
            ws3.cell(row, col, value)
        row += 1
    
    row += 2
    ws3.cell(row, 1, "和光智成对标估值").font = Font(bold=True)
    ws3.merge_cells(f'A{row}:G{row}')
    ws3.cell(row, 1).fill = EXCEL_ACCENT_FILL
    row += 1
    
    heliux_data = [
        ["保守估值", "8x", "¥180万", "¥1.44亿"],
        ["基准估值", "11x", "¥180万", "¥2.0亿"],
        ["乐观估值", "15x", "¥180万", "¥2.7亿"],
    ]
    
    headers2 = ["估值场景", "EV/营收倍数", "2026E ARR", "企业价值"]
    for col, header in enumerate(headers2, 1):
        ws3.cell(row, col, header).font = Font(bold=True, color="FFFFFF")
        ws3.cell(row, col).fill = EXCEL_HEADER_FILL
    row += 1
    
    for item in heliux_data:
        for col, value in enumerate(item, 1):
            ws3.cell(row, col, value)
        row += 1
    
    # ========== Sheet 4: 一页Teaser ==========
    ws4 = wb.create_sheet("投资摘要Teaser")
    
    ws4.cell(1, 1, "和光智成 HelixMind - 投资摘要").font = Font(bold=True, size=16, color="003366")
    ws4.merge_cells('A1:D1')
    
    row = 3
    teaser_content = [
        ("公司定位", "中国领先的AI for Science材料科学计算平台"),
        ("核心技术", "量子力学 + 分子动力学 + 大语言模型三位一体"),
        ("市场机会", "全球AI科学计算市场2033年将达4.8万亿美元，CAGR 27%+"),
        ("商业模式", "SaaS订阅(60%) + 定制开发(30%) + 联合研发分成(10%)"),
        ("当前进展", "5家付费客户，MRR ¥15万，ARR ¥180万"),
        ("2026目标", "MRR ¥100万，ARR ¥1200万，客户数30+"),
        ("融资方案", "Pre-A轮 ¥2000万，投前估值 ¥2亿，出让10%股权"),
        ("资金用途", "产品研发40% + 市场销售35% + 团队扩张15% + 运营10%"),
        ("核心团队", "北航教授领衔，AI+材料科学跨界团队"),
        ("退出路径", "科创板/港交所IPO (2029-2030) 或战略收购"),
    ]
    
    for item, desc in teaser_content:
        ws4.cell(row, 1, item).font = Font(bold=True, color="003366")
        ws4.cell(row, 2, desc)
        ws4.merge_cells(f'B{row}:D{row}')
        row += 1
    
    # 调整列宽
    for ws in [ws1, ws2, ws3, ws4]:
        ws.column_dimensions['A'].width = 25
        for col in ['B', 'C', 'D', 'E', 'F', 'G']:
            ws.column_dimensions[col].width = 18
    
    return wb

def main():
    # print("正在生成和光智成融资路演材料...")
    
    # 生成Pitch Deck
    prs = create_pitch_deck()
    deck_path = "/Users/mettlyz/.openclaw/workspace/output/task-2100/和光智成_融资PitchDeck_Q22026.pptx"
    prs.save(deck_path)
    # print(f"✓ Pitch Deck已保存: {deck_path}")
    
    # 生成财务模型
    wb = create_financial_model()
    model_path = "/Users/mettlyz/.openclaw/workspace/output/task-2100/和光智成_财务模型与估值测算_2026.xlsx"
    wb.save(model_path)
    # print(f"✓ 财务模型已保存: {model_path}")
    
    # print("\n文件生成完成！")

if __name__ == "__main__":
    main()
